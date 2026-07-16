from __future__ import annotations

from importlib.resources import path
import logging
import unicodedata
from concurrent.futures import ProcessPoolExecutor, as_completed
from pathlib import Path

from config import CONFIG, Config

log = logging.getLogger("ingest")
logging.basicConfig(level=logging.INFO, format="%(levelname)s %(name)s: %(message)s")

SUPPORTED = {".pdf", ".pptx", ".docx"}


def clean_text(text: str) -> str:
    text = unicodedata.normalize("NFKC", text)
    text = "".join(ch if (ch.isprintable() or ch in "\n\t") else " "
                   for ch in text)
    out, blank = [], 0
    for ln in (l.strip() for l in text.splitlines()):
        if not ln:
            blank += 1
            if blank <= 1:
                out.append("")
        else:
            blank = 0
            out.append(" ".join(ln.split()))
    return "\n".join(out).strip()


# per-file parse
def _parse_pdf_pymupdf4llm(path: Path) -> list[dict]:
    import pymupdf4llm
    pages = []
    for i, p in enumerate(pymupdf4llm.to_markdown(str(path),
                                                   page_chunks=True), 1):
        t = clean_text(p.get("text", ""))
        if len(t) >= 20:
            pages.append({"source": path.name, "page": i, "text": t})
    return pages


def _parse_pdf_pymupdf(path: Path) -> list[dict]:
    import pymupdf as fitz          
    pages = []
    doc = fitz.open(str(path))
    for i, page in enumerate(doc, 1):
        t = clean_text(page.get_text("text") or "")
        if len(t) >= 20:
            pages.append({"source": path.name, "page": i, "text": t})
    doc.close()
    return pages


def _parse_pdf_pypdf(path: Path) -> list[dict]:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    pages = []
    for i, page in enumerate(reader.pages, 1):
        t = clean_text(page.extract_text() or "")
        if len(t) >= 20:
            pages.append({"source": path.name, "page": i, "text": t})
    return pages


def _parse_pdf_pdfminer(path: Path) -> list[dict]:
    from pdfminer.high_level import extract_pages
    from pdfminer.layout import LTTextContainer
    pages = []
    for i, layout in enumerate(extract_pages(str(path)), 1):
        parts = [el.get_text() for el in layout
                 if isinstance(el, LTTextContainer)]
        t = clean_text("".join(parts))
        if len(t) >= 20:
            pages.append({"source": path.name, "page": i, "text": t})
    return pages


# Parser chain: richest → fastest → pure-Python A → pure-Python B
_PDF_PARSERS: list[tuple[str, object]] = [
    ("pymupdf4llm (markdown+headings)", _parse_pdf_pymupdf4llm),
    ("PyMuPDF plain-text",              _parse_pdf_pymupdf),
    ("pypdf (pure-Python)",             _parse_pdf_pypdf),
    ("pdfminer.six (pure-Python)",      _parse_pdf_pdfminer),
]


def parse_pdf(path: Path) -> list[dict]:
    last_exc: Exception | None = None
    for name, fn in _PDF_PARSERS:
        try:
            pages = fn(path)      
            if pages:
                log.debug("%s: parsed with %s (%d pages)",
                          path.name, name, len(pages))
                return pages
        except ImportError:
            continue
        except Exception as exc:
            last_exc = exc
            log.debug("%s: %s failed (%s)", path.name, name, exc)
            continue
    raise RuntimeError(
        f"All PDF parsers failed for {path.name}. "
        f"Last error: {last_exc}\n"
        f"Fix options:\n"
        f"Install official Python from python.org (not Chocolatey) and "
        f"re-run: pip install PyMuPDF pymupdf4llm\n"
        f"Or install pure-Python fallbacks: pip install pypdf pdfminer.six"
    )


def parse_pptx(path: Path) -> list[dict]:
    from pptx import Presentation
    pages: list[dict] = []
    for i, slide in enumerate(Presentation(path).slides, 1):
        parts = []
        for sh in slide.shapes:
            if sh.has_text_frame:
                parts.append(sh.text_frame.text)
            if getattr(sh, "has_table", False) and sh.has_table:
                for row in sh.table.rows:
                    parts.append(" | ".join(c.text for c in row.cells))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            n = slide.notes_slide.notes_text_frame.text
            if n.strip():
                parts.append("[Speaker notes] " + n)
        t = clean_text("\n".join(parts))
        if t:
            pages.append({"source": path.name, "page": i, "text": t})
    return pages

def parse_docx(path: Path) -> list[dict]:
    from docx import Document
    doc = Document(str(path))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(c.text for c in row.cells))
    t = clean_text("\n".join(parts))
    return [{"source": path.name, "page": 1, "text": t}] if t else []


def _parse_one(path_str: str) -> tuple[str, list[dict], str]:
    path = Path(path_str)
    try:
        suf = path.suffix.lower()
        pages = (parse_pdf(path) if suf == ".pdf"
         else parse_docx(path) if suf == ".docx"
         else parse_pptx(path))
        return path.name, pages, ""
    except Exception as exc:
        return path.name, [], f"{type(exc).__name__}: {exc}"


# incremental run
def ingest(cfg: Config = CONFIG, workers: int | None = None) -> dict:
    from store import Store
    from chunk import build_parent_and_children

    store = Store(cfg)
    files = sorted(p for p in cfg.raw_dir.rglob("*")
                   if p.is_file() and p.suffix.lower() in SUPPORTED)
    if not files:
        raise FileNotFoundError(
            f"No PDF/PPTX under {cfg.raw_dir} - drop your module materials "
            f"there (subfolders are fine).")

    for gone in store.stale_files({p.name for p in files}):
        log.info("Removing deleted file from index: %s", gone)
        store.remove_file(gone)

    todo = [p for p in files if store.file_changed(p)]
    log.info("%d files total, %d new/changed to (re)ingest.", len(files),
             len(todo))
    stats = {"files": len(files), "reingested": 0, "failed": 0, "chunks": 0}

    if not todo:
        return stats

    with ProcessPoolExecutor(max_workers=workers) as pool:
        futures = {pool.submit(_parse_one, str(p)): p for p in todo}
        for fut in as_completed(futures):
            path = futures[fut]
            name, pages, err = fut.result()
            if err:
                log.error("SKIPPING %s: %s", name, err)
                stats["failed"] += 1
                continue
            if not pages:
                log.warning("%s yielded no text - skipped.", name)
                stats["failed"] += 1
                continue
            parents, children = build_parent_and_children(pages, cfg)
            store.replace_file(path, parents, children)
            stats["reingested"] += 1
            stats["chunks"] += len(children)
            log.info("Ingested %-45s %4d pages -> %4d chunks",
                     name, len(pages), len(children))

    log.info("Corpus now holds %d chunks.", store.n_chunks())
    store.close()
    return stats


if __name__ == "__main__":
    print(ingest())
